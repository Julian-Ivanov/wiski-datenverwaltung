from azure.storage.blob import BlobServiceClient
import hashlib
import json
import streamlit as st
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
import re
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from io import BytesIO

# Azure Blob Storage Configuration
connection_string = st.secrets["AZURE_BLOB_CONNECTION_STRING"]
container_name = st.secrets["AZURE_BLOB_CONTAINER_NAME"]
blob_service_client = BlobServiceClient.from_connection_string(connection_string)
container_client = blob_service_client.get_container_client(container_name)

# Azure Search Configuration
search_service_name = st.secrets["AZURE_SEARCH_SERVICE_NAME"]
index_name = st.secrets["AZURE_SEARCH_INDEX_NAME"]
admin_key = st.secrets["AZURE_SEARCH_ADMIN_KEY"]
endpoint = f"https://{search_service_name}.search.windows.net"
search_client = SearchClient(endpoint, index_name, AzureKeyCredential(admin_key))

# Metadata File for Tracking Last Modified Times
metadata_file = "./blob_metadata.json"


def load_previous_metadata():
    """Load previous metadata (last modified times) from a JSON file."""
    if os.path.exists(metadata_file):
        with open(metadata_file, "r") as f:
            return json.load(f)
    return {}


def save_metadata(metadata):
    """Save current metadata (last modified times) to a JSON file."""
    with open(metadata_file, "w") as f:
        json.dump(metadata, f)


def sanitize_key(file_name):
    """Sanitize the file name to generate a valid document key."""
    sanitized = re.sub(r"[^\w\-=]", "_", file_name)
    return sanitized[:1024]  # Truncate to Azure's limit


def extract_text_from_pdf(content):
    """Extract text content from PDF binary data."""
    try:
        # Use PyPDF2 to read the PDF content
        from io import BytesIO
        pdf = PdfReader(BytesIO(content))
        text = ""
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return text
    except Exception as e:
        print(f"Failed to extract text from PDF: {e}")
        return None


def extract_text(file_content, file_extension):
    """Extract text content from a file based on its extension."""
    try:
        if file_extension == "pdf":
            pdf = PdfReader(BytesIO(file_content))
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text
            return text
        elif file_extension == "docx":
            doc = Document(BytesIO(file_content))
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == "pptx":
            presentation = Presentation(BytesIO(file_content))
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif file_extension == "xlsx":
            workbook = openpyxl.load_workbook(BytesIO(file_content))
            text = ""
            for sheet in workbook.sheetnames:
                ws = workbook[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
            return text
        elif file_extension in ["json", "html", "xml", "txt"]:
            return file_content.decode("utf-8")
        else:
            print(f"Unsupported file type: {file_extension}")
            return None
    except Exception as e:
        print(f"Error processing {file_extension} file: {e}")
        return None


def upload_files_to_index(new_files):
    """Upload new or updated blobs to Azure Search index."""
    for blob_name, blob_content in new_files.items():
        file_extension = blob_name.split(".")[-1].lower()
        document_key = sanitize_key(blob_name)

        # Extract content based on file type
        content = extract_text(blob_content, file_extension)
        if content is None:
            print(f"Skipping {blob_name} due to unsupported format or extraction failure.")
            continue

        # Prepare document for indexing
        document = {
            "id": document_key,
            "content": content,
        }

        # Upload to Azure Search
        try:
            results = search_client.upload_documents(documents=[document])
            for result in results:
                if result.succeeded:
                    print(f"Uploaded {blob_name} successfully with ID: {document_key}")
                else:
                    print(f"Failed to upload {blob_name} with ID: {document_key}, Error: {result.error_message}")
        except Exception as e:
            print(f"Failed to upload {blob_name}: {e}")

def delete_removed_files(indexed_keys, current_keys):
    """Delete documents from the index that are no longer in the blob container."""
    keys_to_delete = set(indexed_keys) - set(current_keys)
    if keys_to_delete:
        delete_actions = [{"@search.action": "delete", "id": key} for key in keys_to_delete]
        try:
            results = search_client.upload_documents(documents=delete_actions)
            for result in results:
                if result.succeeded:
                    print(f"Deleted document with ID: {result.key}")
                else:
                    print(f"Failed to delete document with ID: {result.key}, Error: {result.error_message}")
        except Exception as e:
            print(f"Failed to delete documents: {e}")


def main():
    """Main function to synchronize the Blob Storage container with the Azure index."""
    previous_metadata = load_previous_metadata()
    current_metadata = {}
    new_files = {}
    current_file_keys = []

    # Fetch blobs from the container
    for blob in container_client.list_blobs():
        blob_name = blob.name
        last_modified = blob.last_modified.isoformat()  # Use ISO format for serialization
        blob_client = container_client.get_blob_client(blob_name)
        current_file_keys.append(sanitize_key(blob_name))

        # Check if the blob is new or modified
        if blob_name not in previous_metadata or previous_metadata[blob_name] != last_modified:
            blob_content = blob_client.download_blob().readall()
            new_files[blob_name] = blob_content
            current_metadata[blob_name] = last_modified

    # Upload new or updated blobs to the index
    if new_files:
        print(f"New or updated files found: {list(new_files.keys())}")
        upload_files_to_index(new_files)
    else:
        print("No new or updated files to index.")

    # Get keys currently indexed in Azure Search
    indexed_keys = [doc["id"] for doc in search_client.search("")]

    # Delete blobs no longer in the container
    delete_removed_files(indexed_keys, current_file_keys)

    # Save the updated metadata
    save_metadata(current_metadata)


if __name__ == "__main__":
    main()

